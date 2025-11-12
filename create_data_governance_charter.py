#!/usr/bin/env python3
"""
Data Governance Charter Presentation Generator
Creates a comprehensive PowerPoint presentation for an Indian Life Insurance Provider
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide to the presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for item in content_items:
        if isinstance(item, tuple):
            # Tuple format: (text, level)
            p = text_frame.add_paragraph()
            p.text = item[0]
            p.level = item[1]
        else:
            # Simple string
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Remove default placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add left text box
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf_left = left.text_frame
    tf_left.word_wrap = True
    for item in left_content:
        p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
    
    # Add right text box
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    tf_right = right.text_frame
    tf_right.word_wrap = True
    for item in right_content:
        p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
    
    return slide

def create_data_governance_charter_presentation():
    """Create the complete Data Governance Charter presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Data Governance Charter",
        "For Indian Life Insurance Provider\nEstablishing Excellence in Data Management"
    )
    
    # Slide 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        [
            "Purpose: Establish a comprehensive framework for managing data assets",
            "Scope: All data across the organization including customer, policy, claims, and regulatory data",
            "Compliance: Aligned with IRDAI regulations and data protection laws",
            "Objectives:",
            ("Ensure data quality, accuracy, and consistency", 1),
            ("Protect sensitive customer information", 1),
            ("Enable data-driven decision making", 1),
            ("Maintain regulatory compliance", 1),
            ("Establish clear accountability for data assets", 1)
        ]
    )
    
    # Slide 3: Vision and Mission
    add_content_slide(
        prs,
        "Data Governance Vision & Mission",
        [
            "Vision:",
            ("To be recognized as a data-driven organization that leverages trusted, high-quality data to deliver exceptional customer value and maintain regulatory excellence", 1),
            "",
            "Mission:",
            ("Implement robust data governance practices across all business functions", 1),
            ("Ensure data is accurate, accessible, and secure", 1),
            ("Foster a data-driven culture throughout the organization", 1),
            ("Maintain compliance with IRDAI and data protection regulations", 1),
            ("Enable strategic decision-making through trusted data assets", 1)
        ]
    )
    
    # Slide 4: Data Governance Framework
    add_content_slide(
        prs,
        "Data Governance Framework",
        [
            "Core Components:",
            ("Data Governance Council - Strategic oversight and decision-making", 1),
            ("Data Governance Office - Day-to-day operations and coordination", 1),
            ("Subject Area Data Owners - Accountability for specific domains", 1),
            ("Data Stewards - Operational data quality management", 1),
            ("Data Custodians - Technical implementation and security", 1),
            "",
            "Supporting Elements:",
            ("Policies and Standards", 1),
            ("Data Quality Framework", 1),
            ("Security and Privacy Controls", 1),
            ("Metrics and Monitoring", 1)
        ]
    )
    
    # Slide 5: Organizational Structure - Governance Council
    add_content_slide(
        prs,
        "Organizational Structure: Data Governance Council",
        [
            "Chair: Chief Data Officer (CDO) / Chief Information Officer (CIO)",
            "",
            "Members:",
            ("Chief Financial Officer (CFO)", 1),
            ("Chief Risk Officer (CRO)", 1),
            ("Chief Compliance Officer", 1),
            ("Head of Actuarial", 1),
            ("Head of Operations", 1),
            ("Head of IT/Technology", 1),
            ("Head of Customer Service", 1),
            "",
            "Responsibilities:",
            ("Approve data governance policies and standards", 1),
            ("Resolve escalated data-related issues", 1),
            ("Allocate resources for data initiatives", 1),
            ("Review governance metrics quarterly", 1)
        ]
    )
    
    # Slide 6: Subject Area 1 - Policy Administration Data
    add_content_slide(
        prs,
        "Subject Area: Policy Administration Data",
        [
            "Data Owner: Head of Policy Administration / Chief Operations Officer",
            "",
            "Data Stewards:",
            ("Policy Operations Manager", 1),
            ("Premium Administration Lead", 1),
            ("Policy Services Manager", 1),
            "",
            "Key Data Elements:",
            ("Policy details (number, type, status, sum assured)", 1),
            ("Premium information (amount, frequency, payment status)", 1),
            ("Policy riders and add-ons", 1),
            ("Policy lifecycle events (issuance, renewals, lapses, surrenders)", 1),
            ("Beneficiary information", 1),
            ("Maturity and surrender values", 1),
            "",
            "Compliance Requirements: IRDAI guidelines on policy administration"
        ]
    )
    
    # Slide 7: Subject Area 2 - Claims Management Data
    add_content_slide(
        prs,
        "Subject Area: Claims Management Data",
        [
            "Data Owner: Head of Claims / Chief Claims Officer",
            "",
            "Data Stewards:",
            ("Claims Operations Manager", 1),
            ("Claims Underwriting Lead", 1),
            ("Claims Settlement Manager", 1),
            "",
            "Key Data Elements:",
            ("Claim registration details", 1),
            ("Claimant information and documentation", 1),
            ("Claim assessment and investigation records", 1),
            ("Settlement amounts and payment details", 1),
            ("Claim rejection reasons and appeals", 1),
            ("Fraud detection indicators", 1),
            "",
            "Compliance Requirements: IRDAI claim settlement regulations"
        ]
    )
    
    # Slide 8: Subject Area 3 - Customer/Policyholder Data
    add_content_slide(
        prs,
        "Subject Area: Customer/Policyholder Data",
        [
            "Data Owner: Head of Customer Service / Chief Customer Officer",
            "",
            "Data Stewards:",
            ("Customer Experience Manager", 1),
            ("CRM System Administrator", 1),
            ("Customer Data Quality Analyst", 1),
            "",
            "Key Data Elements:",
            ("Personal information (name, DOB, contact details, PAN, Aadhaar)", 1),
            ("KYC documentation and verification status", 1),
            ("Customer preferences and communication history", 1),
            ("Relationship mapping (family members, nominees)", 1),
            ("Customer segmentation and profiling data", 1),
            ("Complaints and grievances", 1),
            "",
            "Compliance Requirements: DPDP Act 2023, IRDAI customer protection guidelines"
        ]
    )
    
    # Slide 9: Subject Area 4 - Agent/Distribution Data
    add_content_slide(
        prs,
        "Subject Area: Agent/Distribution Data",
        [
            "Data Owner: Head of Distribution / Chief Distribution Officer",
            "",
            "Data Stewards:",
            ("Agency Operations Manager", 1),
            ("Bancassurance Partnerships Lead", 1),
            ("Distribution Analytics Manager", 1),
            "",
            "Key Data Elements:",
            ("Agent registration and licensing information", 1),
            ("Agent performance metrics and commissions", 1),
            ("Training and certification records", 1),
            ("Distribution channel data (branch, online, bancassurance)", 1),
            ("Partner/intermediary agreements", 1),
            ("Agent hierarchy and reporting structure", 1),
            "",
            "Compliance Requirements: IRDAI agent licensing norms"
        ]
    )
    
    # Slide 10: Subject Area 5 - Actuarial and Risk Data
    add_content_slide(
        prs,
        "Subject Area: Actuarial and Risk Data",
        [
            "Data Owner: Appointed Actuary / Head of Actuarial",
            "",
            "Data Stewards:",
            ("Actuarial Analyst Lead", 1),
            ("Risk Modeling Specialist", 1),
            ("Reserving Manager", 1),
            "",
            "Key Data Elements:",
            ("Mortality and morbidity experience data", 1),
            ("Actuarial assumptions and parameters", 1),
            ("Reserve calculations and valuations", 1),
            ("Pricing models and rate tables", 1),
            ("Reinsurance treaty data", 1),
            ("Risk assessment and underwriting guidelines", 1),
            "",
            "Compliance Requirements: IRDAI actuarial reporting requirements"
        ]
    )
    
    # Slide 11: Subject Area 6 - Financial and Accounting Data
    add_content_slide(
        prs,
        "Subject Area: Financial and Accounting Data",
        [
            "Data Owner: Chief Financial Officer (CFO)",
            "",
            "Data Stewards:",
            ("Financial Controller", 1),
            ("Investment Accounting Manager", 1),
            ("GL Accounting Lead", 1),
            "",
            "Key Data Elements:",
            ("General ledger and chart of accounts", 1),
            ("Premium collection and revenue recognition", 1),
            ("Investment portfolio data", 1),
            ("Expense allocation and cost centers", 1),
            ("Financial statements and reports", 1),
            ("Budget and forecast data", 1),
            "",
            "Compliance Requirements: Ind AS, IRDAI financial reporting norms"
        ]
    )
    
    # Slide 12: Subject Area 7 - Regulatory and Compliance Data
    add_content_slide(
        prs,
        "Subject Area: Regulatory and Compliance Data",
        [
            "Data Owner: Chief Compliance Officer",
            "",
            "Data Stewards:",
            ("Regulatory Reporting Manager", 1),
            ("Compliance Monitoring Lead", 1),
            ("Audit and Control Specialist", 1),
            "",
            "Key Data Elements:",
            ("IRDAI regulatory returns and filings", 1),
            ("AML/CFT transaction monitoring data", 1),
            ("Solvency and capital adequacy metrics", 1),
            ("Internal audit findings and remediation", 1),
            ("Policy on data localization and privacy", 1),
            ("Regulatory correspondence and approvals", 1),
            "",
            "Compliance Requirements: IRDAI regulations, PMLA, DPDP Act"
        ]
    )
    
    # Slide 13: Subject Area 8 - Underwriting Data
    add_content_slide(
        prs,
        "Subject Area: Underwriting Data",
        [
            "Data Owner: Head of Underwriting / Chief Underwriting Officer",
            "",
            "Data Stewards:",
            ("Underwriting Operations Manager", 1),
            ("Medical Underwriting Lead", 1),
            ("Risk Selection Analyst", 1),
            "",
            "Key Data Elements:",
            ("Proposal and application data", 1),
            ("Medical examination results and reports", 1),
            ("Risk assessment scores and decisions", 1),
            ("Underwriting guidelines and rules", 1),
            ("Loading and exclusion data", 1),
            ("Declined and postponed case records", 1),
            "",
            "Compliance Requirements: IRDAI underwriting norms, medical confidentiality"
        ]
    )
    
    # Slide 14: Roles and Responsibilities - Data Owners
    add_content_slide(
        prs,
        "Roles and Responsibilities: Data Owners",
        [
            "Accountability for specific subject areas of data",
            "",
            "Key Responsibilities:",
            ("Define data requirements and business rules for their domain", 1),
            ("Approve access rights and data sharing agreements", 1),
            ("Ensure compliance with relevant regulations", 1),
            ("Resolve data-related issues within their domain", 1),
            ("Participate in Data Governance Council meetings", 1),
            ("Approve data quality standards for their subject area", 1),
            ("Sponsor data quality improvement initiatives", 1),
            ("Review and approve major changes to data structures", 1),
            "",
            "Authority:",
            ("Final decision-making for their subject area", 1),
            ("Budget allocation for data quality initiatives", 1)
        ]
    )
    
    # Slide 15: Roles and Responsibilities - Data Stewards
    add_content_slide(
        prs,
        "Roles and Responsibilities: Data Stewards",
        [
            "Day-to-day operational management of data quality",
            "",
            "Key Responsibilities:",
            ("Monitor data quality metrics and KPIs", 1),
            ("Identify and resolve data quality issues", 1),
            ("Maintain data definitions and business glossary", 1),
            ("Coordinate with IT for technical implementations", 1),
            ("Conduct data quality audits", 1),
            ("Train business users on data standards", 1),
            ("Document data lineage and transformations", 1),
            ("Manage metadata and data catalogs", 1),
            "",
            "Reporting:",
            ("Report to Data Owners on quality metrics", 1),
            ("Escalate critical issues to Data Governance Office", 1)
        ]
    )
    
    # Slide 16: Roles and Responsibilities - Data Custodians
    add_content_slide(
        prs,
        "Roles and Responsibilities: Data Custodians",
        [
            "IT teams responsible for technical implementation and security",
            "",
            "Key Responsibilities:",
            ("Implement technical controls for data security", 1),
            ("Manage database systems and data storage", 1),
            ("Execute backup and recovery procedures", 1),
            ("Implement access controls and authentication", 1),
            ("Monitor system performance and availability", 1),
            ("Deploy data quality tools and technologies", 1),
            ("Maintain data integration and ETL processes", 1),
            ("Ensure data encryption and protection", 1),
            "",
            "Coordination:",
            ("Work with Data Stewards to implement controls", 1),
            ("Support Data Owners with technical requirements", 1)
        ]
    )
    
    # Slide 17: Data Governance Policies and Standards
    add_content_slide(
        prs,
        "Data Governance Policies and Standards",
        [
            "Core Policies:",
            ("Data Classification and Handling Policy", 1),
            ("Data Access and Authorization Policy", 1),
            ("Data Retention and Archival Policy", 1),
            ("Data Privacy and Protection Policy", 1),
            ("Data Quality Management Policy", 1),
            ("Master Data Management Policy", 1),
            ("Data Sharing and Third-Party Management Policy", 1),
            "",
            "Standards:",
            ("Data naming conventions and standards", 1),
            ("Data definition and metadata standards", 1),
            ("Data quality thresholds and SLAs", 1),
            ("Documentation and change management standards", 1)
        ]
    )
    
    # Slide 18: Data Quality Management
    add_content_slide(
        prs,
        "Data Quality Management Framework",
        [
            "Data Quality Dimensions:",
            ("Accuracy - Data correctly represents real-world values", 1),
            ("Completeness - All required data is present", 1),
            ("Consistency - Data is uniform across systems", 1),
            ("Timeliness - Data is available when needed", 1),
            ("Validity - Data conforms to business rules", 1),
            ("Uniqueness - No duplicate records exist", 1),
            "",
            "Quality Management Processes:",
            ("Data profiling and assessment", 1),
            ("Issue identification and root cause analysis", 1),
            ("Remediation and improvement initiatives", 1),
            ("Ongoing monitoring and reporting", 1),
            ("Continuous improvement programs", 1)
        ]
    )
    
    # Slide 19: Data Security and Privacy
    add_content_slide(
        prs,
        "Data Security and Privacy",
        [
            "Security Measures:",
            ("Role-based access control (RBAC)", 1),
            ("Data encryption at rest and in transit", 1),
            ("Network security and firewalls", 1),
            ("Regular security audits and penetration testing", 1),
            ("Incident response and breach management", 1),
            "",
            "Privacy Compliance:",
            ("DPDP Act 2023 compliance framework", 1),
            ("Consent management for customer data", 1),
            ("Right to access, correction, and erasure", 1),
            ("Data localization requirements", 1),
            ("Privacy impact assessments", 1),
            ("IRDAI data protection guidelines adherence", 1)
        ]
    )
    
    # Slide 20: Implementation Roadmap
    add_content_slide(
        prs,
        "Implementation Roadmap",
        [
            "Phase 1 (Months 1-3): Foundation",
            ("Establish Data Governance Council and Office", 1),
            ("Appoint Data Owners and Stewards", 1),
            ("Define initial policies and standards", 1),
            ("Conduct current state assessment", 1),
            "",
            "Phase 2 (Months 4-6): Core Implementation",
            ("Implement data quality framework", 1),
            ("Deploy metadata management tools", 1),
            ("Establish data quality metrics and monitoring", 1),
            ("Launch training and awareness programs", 1),
            "",
            "Phase 3 (Months 7-12): Expansion",
            ("Roll out across all subject areas", 1),
            ("Implement advanced analytics and reporting", 1),
            ("Mature data quality processes", 1),
            ("Continuous improvement and optimization", 1)
        ]
    )
    
    # Slide 21: Metrics and KPIs
    add_content_slide(
        prs,
        "Data Governance Metrics and KPIs",
        [
            "Data Quality Metrics:",
            ("% of critical data elements meeting quality thresholds", 1),
            ("Number of data quality issues identified and resolved", 1),
            ("Data completeness scores by subject area", 1),
            ("Data accuracy rates for key business processes", 1),
            "",
            "Operational Metrics:",
            ("% of data with assigned owners and stewards", 1),
            ("Average time to resolve data issues", 1),
            ("User satisfaction with data quality", 1),
            ("% of employees trained on data governance", 1),
            "",
            "Compliance Metrics:",
            ("% compliance with IRDAI reporting timelines", 1),
            ("Number of data-related regulatory findings", 1),
            ("% of policies reviewed and updated annually", 1)
        ]
    )
    
    # Slide 22: Benefits and Value Proposition
    add_content_slide(
        prs,
        "Benefits and Value Proposition",
        [
            "Business Benefits:",
            ("Improved decision-making through trusted data", 1),
            ("Enhanced customer experience and satisfaction", 1),
            ("Reduced operational risks and errors", 1),
            ("Faster time-to-market for new products", 1),
            ("Better regulatory compliance and reduced penalties", 1),
            "",
            "Organizational Impact:",
            ("Clear accountability for data assets", 1),
            ("Consistent data definitions across the enterprise", 1),
            ("Improved collaboration between business and IT", 1),
            ("Enhanced data literacy and awareness", 1),
            ("Foundation for advanced analytics and AI initiatives", 1),
            "",
            "Competitive Advantage:",
            ("Data-driven culture and innovation", 1)
        ]
    )
    
    # Slide 23: Governance Review and Continuous Improvement
    add_content_slide(
        prs,
        "Governance Review and Continuous Improvement",
        [
            "Regular Review Cycles:",
            ("Monthly: Data Steward working sessions", 1),
            ("Quarterly: Data Governance Council meetings", 1),
            ("Annually: Comprehensive governance assessment", 1),
            "",
            "Continuous Improvement Focus Areas:",
            ("Policy and standards updates", 1),
            ("Technology and tool enhancements", 1),
            ("Process optimization", 1),
            ("Skills and capability development", 1),
            ("Industry best practice adoption", 1),
            "",
            "Change Management:",
            ("Regular communication and updates", 1),
            ("Stakeholder feedback incorporation", 1),
            ("Lessons learned documentation", 1)
        ]
    )
    
    # Slide 24: Appendix - Key Regulations
    add_content_slide(
        prs,
        "Appendix: Key Regulatory References",
        [
            "IRDAI Regulations:",
            ("IRDAI (Insurance Products) Regulations", 1),
            ("IRDAI (Protection of Policyholders' Interests) Regulations", 1),
            ("IRDAI (Actuarial Report and Abstract) Regulations", 1),
            ("IRDAI (Preparation of Financial Statements) Regulations", 1),
            ("IRDAI (Outsourcing of Activities) Guidelines", 1),
            "",
            "Data Protection and Privacy:",
            ("Digital Personal Data Protection Act (DPDP) 2023", 1),
            ("Information Technology Act, 2000", 1),
            ("Prevention of Money Laundering Act (PMLA) 2002", 1),
            "",
            "Industry Standards:",
            ("ISO 27001 - Information Security Management", 1),
            ("ISO 8000 - Data Quality", 1)
        ]
    )
    
    # Slide 25: Contact and Support
    add_content_slide(
        prs,
        "Contact and Support",
        [
            "Data Governance Office:",
            ("Email: datagovernance@company.com", 1),
            ("Phone: [Contact Number]", 1),
            ("Intranet: [Data Governance Portal URL]", 1),
            "",
            "For Queries:",
            ("Data Access Requests: Contact your Data Owner", 1),
            ("Data Quality Issues: Contact your Data Steward", 1),
            ("Technical Support: Contact Data Custodian/IT Helpdesk", 1),
            ("Policy Questions: Contact Data Governance Office", 1),
            "",
            "Resources:",
            ("Data Governance Portal with policies and procedures", 1),
            ("Training materials and e-learning modules", 1),
            ("Data catalog and business glossary", 1),
            ("Monthly newsletters and updates", 1)
        ]
    )
    
    # Save the presentation
    output_file = "/home/runner/work/test-repo/test-repo/Data_Governance_Charter_Indian_Life_Insurance.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
    return output_file

if __name__ == "__main__":
    create_data_governance_charter_presentation()
